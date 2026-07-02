import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

out = r"D:\gzp\研究生\半导体器件建模与仿真\NN_device_model\ppt\NN紧凑模型文献汇报.pptx"
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ---- Color Palette ----
C_DARK   = RGBColor(0x1a,0x1a,0x2e)
C_ACCENT = RGBColor(0x00,0x6d,0x77)
C_ACC2   = RGBColor(0xe8,0x95,0x22)
C_WHITE  = RGBColor(0xff,0xff,0xff)
C_LIGHT  = RGBColor(0xf0,0xf1,0xf5)
C_GRAY   = RGBColor(0x55,0x55,0x55)
C_SUB    = RGBColor(0xaa,0xaa,0xaa)
C_RED    = RGBColor(0xcc,0x33,0x33)
C_GREEN  = RGBColor(0x2e,0x7d,0x32)
C_BLUE   = RGBColor(0x15,0x65,0xc0)
FONT     = 'Microsoft YaHei'
SPACER   = Inches(0.04)

def slide_bg(s, c=C_DARK):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = c

def rect(s, l, t, w, h, c):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = c
    sh.line.fill.background()
    return sh

def rrect(s, l, t, w, h, c):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = c
    sh.line.fill.background()
    return sh

def tb(s, l, t, w, h, txt, sz=18, bold=False, c=C_WHITE, align=PP_ALIGN.LEFT):
    bx = s.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = txt
    p.font.size = Pt(sz); p.font.bold = bold
    p.font.color.rgb = c; p.font.name = FONT
    p.alignment = align
    return bx

def ml(s, l, t, w, h, lines, sz=16, c=C_WHITE, ls=1.35):
    bx = s.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame; tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text = ln; p.font.size = Pt(sz)
        p.font.color.rgb = c; p.font.name = FONT
        p.space_after = Pt(sz*(ls-1))
    return bx

def hdr_bar(s, title, subtitle=""):
    rect(s, Inches(0), Inches(0), prs.slide_width, Inches(1.2), C_ACCENT)
    tb(s, Inches(0.8), Inches(0.2), Inches(11.5), Inches(0.7), title, sz=30, bold=True)
    if subtitle:
        tb(s, Inches(0.8), Inches(0.85), Inches(11.5), Inches(0.35), subtitle, sz=14, c=C_SUB)

def card(s, l, t, w, h, title, bullets, tsz=18, bsz=14):
    rrect(s, l, t, w, h, C_LIGHT)
    y1 = Emu(t.emu + Inches(0.15).emu)
    tb(s, Emu(l.emu + Inches(0.25).emu), y1, Emu(w.emu - Inches(0.5).emu), Inches(0.4), title, sz=tsz, bold=True, c=C_DARK)
    y2 = Emu(t.emu + Inches(0.6).emu)
    lbs = ["  "+b for b in bullets]
    ml(s, Emu(l.emu + Inches(0.25).emu), y2, Emu(w.emu - Inches(0.5).emu), Emu(h.emu - Inches(0.7).emu), lbs, sz=bsz, c=C_GRAY)

def notes(nf, txt):
    nf.add_paragraph()
    nf.paragraphs[0].text = txt

# ===== S1: Title =====
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(s)
rect(s, Inches(0), Inches(0), prs.slide_width, Inches(0.06), C_ACC2)
tb(s, Inches(1.5), Inches(1.6), Inches(10), Inches(1.2), "神经网络紧凑模型文献汇报", sz=42, bold=True)
tb(s, Inches(1.5), Inches(2.9), Inches(10), Inches(0.6), "A SPICE-Compatible Neural Network Compact Model for Efficient IC Simulations", sz=20, c=C_ACCENT)
rect(s, Inches(1.5), Inches(3.6), Inches(2.8), Inches(0.04), C_ACC2)
tb(s, Inches(1.5), Inches(3.85), Inches(10), Inches(0.5), "Tung C.T., Salahuddin S., Hu C.  |  UC Berkeley  |  2024 IEEE", sz=16, c=C_SUB)
tb(s, Inches(1.5), Inches(4.3), Inches(10), Inches(0.5), "研究方向：NN/ML驱动的半导体器件紧凑模型与SPICE兼容电路仿真", sz=15, c=C_SUB)
tb(s, Inches(1.5), Inches(4.8), Inches(10), Inches(0.5), "延伸阅读：BSIM-NN Framework (TED 2023) / BSIM-NN ML (2025)", sz=14, c=C_SUB)
notes(s.notes_slide.notes_text_frame, "各位老师同学大家好，今天我汇报的文献是UC Berkeley胡正明院士课题组发表在2024年IEEE上的论文《A SPICE-Compatible Neural Network Compact Model for Efficient IC Simulations》。该工作首次系统性地提出了一种将神经网络直接嵌入SPICE仿真引擎的紧凑模型框架，并通过Verilog-A实现和电路级基准测试验证了其可行性和仿真加速能力。同时我会结合该组2023年发表在TED上的BSIM-NN框架论文和2025年的BSIM-NN机器学习紧凑模型论文，对NN紧凑模型的完整技术链条进行梳理。")

# ===== S2: 研究背景 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(s)
hdr_bar(s, "晶体管紧凑模型是IC设计的基石，传统物理模型面临建模效率瓶颈", "#02 研究背景")
card(s, Inches(0.5), Inches(1.5), Inches(6), Inches(2.6), "紧凑模型的核心角色", [
    "晶体管模型是电路仿真中 KCL/KVL 矩阵方程的基础构件",
    "需提供终端电流 I(V) 和电荷 Q(V) 的非线性解析表达式及其导数",
    "从 BSIM3/4 到 BSIM-CMG，参数数量随工艺节点指数增长（>2000）",
    "模型开发周期长：新器件（FinFET/GAA）从物理建模到工业可用需数年",
], tsz=18, bsz=15)
card(s, Inches(6.8), Inches(1.5), Inches(6), Inches(2.6), "现有方法的不足", [
    "物理紧凑模型：解析方程推导复杂，需大量近似（GCA、CSA等）",
    "数据表模型：历史久但电路类型受限，未成为主流工业标准",
    "现有NN尝试：仅覆盖IV或有限工况，缺乏完整的电路级仿真验证",
    "AI/ML模型缺乏工业级基准测试：噪声、统计变异、Gummel对称性等",
], tsz=18, bsz=15)
card(s, Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.6), "关键趋势：DTCO范式下的模型加速需求", [
    "设计-工艺协同优化（DTCO）成为技术发展新范式，要求模型周转时间从月缩短到天甚至小时",
    "AI/ML方法有望将模型开发从"物理推导-参数提取"两阶段流程转变为数据驱动的端到端流程",
    "SPICE兼容性是NN模型能否"落地"的第一道门槛：需保证导数连续性、Jacobian非奇异、NR收敛",
], tsz=18, bsz=15)
notes(s.notes_slide.notes_text_frame, "紧凑模型是整个EDA生态的基石。传统物理模型虽然精度高，但开发周期长、参数提取复杂，尤其对于新兴器件难以迅速响应。近年来设计工艺协同优化DTCO成为新范式，要求模型能快速迭代。AI/ML方法正是看中了这一点。但要真正落地，SPICE兼容性是第一道门槛——模型必须保证导数连续、Jacobian矩阵非奇异、Newton-Raphson迭代收敛。目前已有的NN模型尝试大多只覆盖了IV特性，缺乏完整的电路级闭环验证。")

# ===== S3: 科学问题 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(s)
hdr_bar(s, "能否用神经网络替代物理方程，构建兼具精度与速度的SPICE兼容紧凑模型？", "#03 科学问题与研究动机")
card(s, Inches(0.5), Inches(1.5), Inches(4), Inches(2.5), "核心科学问题", [
    "NN模型能否达到物理紧凑模型的IV/CV拟合精度",
    "NN模型在SPICE引擎中能否保证数值求解的收敛性",
    "NN模型能否在大规模电路中实现仿真加速而不损失精度",
], tsz=18, bsz=15)
card(s, Inches(4.8), Inches(1.5), Inches(4), Inches(2.5), "研究动机（技术驱动）", [
    "先进FET物理日趋复杂，传统建模周期太长",
    "NN可自动从数据学习多维度映射(L,W,EOT)",
    "直接乘法替代Verilog-A循环可显著提升速度",
], tsz=18, bsz=15)
card(s, Inches(9.1), Inches(1.5), Inches(4), Inches(2.5), "研究动机（方法论驱动）", [
    "ISRU激活函数避免了指数计算，适合VA实现",
    "导数辅助损失函数提升边界精度与平滑性",
    "用Python自动生成VA代码，打通训练到仿真链路",
], tsz=18, bsz=15)
rrect(s, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.6), C_DARK)
rect(s, Inches(0.55), Inches(4.35), Inches(0.05), Inches(2.5), C_ACC2)
ml(s, Inches(1.0), Inches(4.6), Inches(11), Inches(2.0), [
    "论文目标定位（Section I, II）",
    "  构建完整的 IV + QV 双网络 NN 紧凑模型，覆盖所有终端电流和电荷，包含几何尺寸依赖性",
    "  研究不同激活函数（Sigmoid / Tanh / ISRU）对电路仿真速度和收敛性的影响",
    "  以直接乘法（非循环）方式实现Verilog-A代码，并在多种电路中基准测试 vs BSIM-CMG",
], sz=16, c=C_WHITE)
notes(s.notes_slide.notes_text_frame, "本文的科学问题可以归纳为三个层面。第一，NN能否在精度上匹敌物理模型？这涉及IV和CV的双重拟合。第二，NN模型能否在SPICE这种对数值稳定性要求极高的求解器中正常工作？这涉及激活函数的选择和导数的光滑性。第三，NN模型能否真的带来仿真加速？这涉及Verilog-A实现策略。作者的动机很明确：先进器件的物理越来越复杂，而NN可以从数据中自动学习多维度映射，关键在于找到一种SPICE兼容的实现方式。")

# ===== S4: 相关工作 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(s)
hdr_bar(s, "NN紧凑模型领域已有初步探索，但在模型完备性和电路级验证方面存在明显gap", "#04 相关工作与现有方法不足")
card(s, Inches(0.5), Inches(1.5), Inches(6), Inches(2.5), "已有NN建模工作", [
    "Li et al. (2016) 物理启发NN：仅覆盖IV，无CV/Gummel对称性验证",
    "Wang et al. (2021) ANN紧凑模型：包含IV+Cgg，仅10神经元，未做大电路测试",
    "Qi et al. (2023) 知识基NN SPICE模型：针对2D材料FET，需物理先验知识",
    "Huang & Wang (2023) 物理基ANN批量生产：依赖物理分区采样策略",
], tsz=17, bsz=14)
card(s, Inches(6.8), Inches(1.5), Inches(6), Inches(2.5), "本论文与其前序工作的延续", [
    "Tung et al. (2022, TED): 首次验证NN IV模型+CV网络，40x提速于小电路",
    "Tung & Hu (2023, TED): 扩展为BSIM-NN框架，IV+QV+噪声+变异，Gummel测试",
    "Tung et al. (2025): 完整BSIM-NN ML：含NQS、自热、RF噪声，481%仿真提速",
    "本文(2024): 聚焦SPICE兼容实现策略，系统对比激活函数和VA编码方式",
], tsz=17, bsz=14)
card(s, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.6), "现有工作的共同不足 (McAndrew et al. 2025基准测试清单)", [
    "多数工作仅使用合成数据（TCAD或模型生成），未在实测数据上验证，对测量噪声和局部变异的鲁棒性未知",
    "缺乏系统性的SPICE模型质量测试：Gummel对称性(4阶导数)、谐波平衡、瞬态电荷守恒等工业标准测试大多缺失",
    "几何/温度/工艺变异等多维度覆盖不足，标度律(Scaling)能力的物理自洽性未得到严格证明",
], tsz=17, bsz=14)
notes(s.notes_slide.notes_text_frame, "在相关工作部分，我们梳理了两条线。一条是更广泛的NN建模尝试，从2016年的物理启发NN到2023年的知识基模型，虽然思路丰富，但要么仅覆盖IV，要么缺乏电路级验证。另一条是UC Berkeley课题组自身的迭代路线——从2022年的概念验证到2025年的完整BSIM-NN，本文2024年这项工作起到了承上启下的关键作用：系统解决了SPICE兼容性这一工程难题。此外，McAndrew等人在2025年提出了AI/ML SPICE模型的工业级基准清单，指出当前领域普遍存在合成数据验证、缺乏工业标准测试等不足。")

# ===== S5: 核心思想 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(s)
hdr_bar(s, "核心思想：用NN替代物理方程的数值计算，但保留紧凑模型在SPICE中的函数接口形式", "#05 核心思想与关键洞察")
# Flow diagram
# Box 1
rrect(s, Inches(0.3), Inches(1.7), Inches(2.8), Inches(1.8), C_ACCENT)
tb(s, Inches(0.5), Inches(1.85), Inches(2.4), Inches(0.4), "NN = 物理方程替代", sz=17, bold=True)
ml(s, Inches(0.5), Inches(2.3), Inches(2.4), Inches(0.9), ["NN替代BSIM中的", "分区域解析方程", "统一拟合全局特性"], sz=13)
# Arrow 1
tb(s, Inches(3.2), Inches(2.3), Inches(0.8), Inches(0.6), "->", sz=28, bold=True, c=C_ACC2, align=PP_ALIGN.CENTER)
# Box 2
rrect(s, Inches(3.9), Inches(1.7), Inches(2.8), Inches(1.8), C_ACCENT)
tb(s, Inches(4.1), Inches(1.85), Inches(2.4), Inches(0.4), "关键输出变换", sz=17, bold=True)
ml(s, Inches(4.1), Inches(2.3), Inches(2.4), Inches(0.9), ["ID = VDS exp(y1)", "保证 VDS=0 时 ID=0", "IG用符号分离+平滑"], sz=13)
# Arrow 2
tb(s, Inches(6.8), Inches(2.3), Inches(0.8), Inches(0.6), "->", sz=28, bold=True, c=C_ACC2, align=PP_ALIGN.CENTER)
# Box 3
rrect(s, Inches(7.5), Inches(1.7), Inches(2.8), Inches(1.8), C_ACCENT)
tb(s, Inches(7.7), Inches(1.85), Inches(2.4), Inches(0.4), "ISRU激活函数", sz=17, bold=True)
ml(s, Inches(7.7), Inches(2.3), Inches(2.4), Inches(0.9), ["f(x)=x/sqrt(1+x^2)", "无指数函数, 导数光滑", "VA实现效率最优"], sz=13)
# Arrow 3
tb(s, Inches(10.4), Inches(2.3), Inches(0.8), Inches(0.6), "->", sz=28, bold=True, c=C_ACC2, align=PP_ALIGN.CENTER)
# Box 4
rrect(s, Inches(11.1), Inches(1.7), Inches(2.0), Inches(1.8), C_ACCENT)
tb(s, Inches(11.25), Inches(1.85), Inches(1.7), Inches(0.4), "直接乘法VA", sz=17, bold=True)
ml(s, Inches(11.25), Inches(2.3), Inches(1.7), Inches(0.9), ["展开矩阵乘法", "无循环/无数组", "~30x速度提升"], sz=13)

# Key insights
card(s, Inches(0.3), Inches(3.8), Inches(6), Inches(3.2), "三个关键洞察", [
    "输出变换 = NN精度保障 + 物理约束：ID=VDS x exp(y1)自动满足VDS=0边界条件；IG符号分离用平滑smoothing函数确保对数变换可学习",
    "导数辅助损失 = SPICE收敛保障：损失函数包含gm, gds, gm', gds'的RMS误差，训练出的NN天然具有光滑导数，适配NR求解器",
    "直接乘法实现 = 仿真速度保障：Verilog-A中展开所有矩阵运算为逐元素乘加，避免循环和数组，在一个17级环振中比循环实现快30倍",
], tsz=17, bsz=14)
card(s, Inches(6.6), Inches(3.8), Inches(6), Inches(3.2), "与物理模型的本质区别", [
    "物理模型：分区建模（亚阈值/线性/饱和）+ 经验光滑函数拼接",
    "NN模型：单网络全局拟合，激活函数本身提供天然光滑性",
    "物理模型参数有物理含义（VTH0, U0...）；NN模型参数为抽象权重",
    "NN模型的"可解释性"来自导数损失对物理量(gm,gds)的监督学习",
], tsz=17, bsz=14)
notes(s.notes_slide.notes_text_frame, "这一页讲解本文的核心思想和技术路线。论文的核心贡献可以分为四个层面：第一，用神经网络替代传统BSIM模型中的分区域解析方程；第二，通过ID=VDS*exp(y1)这种输出变换将物理边界条件嵌入NN结构；第三，选择ISRU激活函数，在无指数计算的前提下保证导数光滑；第四，在Verilog-A实现上将矩阵乘法展开为直接乘法，避免低效的循环结构。最后一点值得强调的是，NN模型的可解释性并非来自权重本身，而是通过导数损失函数对物理量（如gm和gds）进行监督学习来保证的。")

# ===== S6: 方法框架 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(s)
hdr_bar(s, "方法框架：IV+QV双网络架构，ISRU激活 + 导数辅助损失 + 直接乘法Verilog-A", "#06 方法框架")
# Main architecture boxes
rrect(s, Inches(0.3), Inches(1.5), Inches(3.6), Inches(1.8), C_ACCENT)
tb(s, Inches(0.5), Inches(1.6), Inches(3.2), Inches(0.4), "输入层", sz=18, bold=True)
ml(s, Inches(0.5), Inches(2.05), Inches(3.2), Inches(1.1), ["VGS, VDS, HFIN, L, EOT", "(+ T, T_FIN etc.)", "Z-score归一化"], sz=14)
tb(s, Inches(4.0), Inches(2.1), Inches(0.7), Inches(0.6), "->", sz=24, bold=True, c=C_ACC2, align=PP_ALIGN.CENTER)
rrect(s, Inches(4.6), Inches(1.5), Inches(3.6), Inches(1.8), C_ACCENT)
tb(s, Inches(4.8), Inches(1.6), Inches(3.2), Inches(0.4), "2个隐藏层 x 10神经元", sz=18, bold=True)
ml(s, Inches(4.8), Inches(2.05), Inches(3.2), Inches(1.1), ["激活函数对比：", "Sigmoid / Tanh / ISRU", "ISRU最佳(无exp, 更简洁)"], sz=14)
tb(s, Inches(8.3), Inches(2.1), Inches(0.7), Inches(0.6), "->", sz=24, bold=True, c=C_ACC2, align=PP_ALIGN.CENTER)
rrect(s, Inches(8.9), Inches(1.5), Inches(4.2), Inches(1.8), C_ACCENT)
tb(s, Inches(9.1), Inches(1.6), Inches(3.8), Inches(0.4), "输出层 (IV + QV 双网络)", sz=18, bold=True)
ml(s, Inches(9.1), Inches(2.05), Inches(3.8), Inches(1.1), ["IV: y1=ln(ID/VDS), y2p, y2n(IG)", "QV: QG, QS, QD -> Cgg,Cgd,Cgs", "导数损失: gm,gds,gm',gds'"], sz=14)

# Loss function detail card
card(s, Inches(0.3), Inches(3.6), Inches(7.5), Inches(3.3), "损失函数设计 (Eq. 4, Section II)", [
    "L_IV = a RMS(y1) + b RMS(gm) + c RMS(gds) + d RMS(gm') + e RMS(gds') + f RMS(y2p) + f RMS(y2n)",
    "- y1项监督ID精度; gm/gds项保证导数准确性; gm'/gds'二阶导数项保证曲率光滑",
    "- 系数a~f通过网格搜索确定，训练策略：先用简单损失(y1+gm+gds)预训练，再引入高阶导数精调",
], tsz=16, bsz=14)

# Additional details
card(s, Inches(8.1), Inches(3.6), Inches(5), Inches(3.3), "训练与部署", [
    "数据：BSIM-CMG生成 IV/CV 数据（L=14~24nm, HFIN=38~54nm, EOT=0.68~0.88nm）",
    "训练：2层x10神经元，共~200参数",
    "导出：Python自动生成Verilog-A",
    "加速：直接乘法展开（非循环）",
], tsz=16, bsz=14)

notes(s.notes_slide.notes_text_frame, "本页详细介绍方法框架。模型采用IV和QV双网络架构，输入层包含VGS、VDS及器件几何参数HFIN、L、EOT。两个隐藏层各10个神经元，对比了Sigmoid、Tanh和ISRU三种激活函数，其中ISRU因其函数形式简洁且不含指数运算而表现最佳。损失函数的设计是本文方法论的亮点：不仅包含输出量y1的RMS误差，还包含一阶导数gm、gds和二阶导数gm'、gds'的误差项。这确保了训练出的NN具有物理上正确的电导特性，这是SPICE收敛的关键保障。训练完成后，通过Python脚本自动将权重导出为Verilog-A代码，所有矩阵乘法展开为逐元素乘加。")

# ===== S7: 实验设计 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(s)
hdr_bar(s, "实验设计：多维度精度验证 + 跨规模电路基准测试 + 多种激活函数对比", "#07 实验设计")
card(s, Inches(0.3), Inches(1.5), Inches(4), Inches(2.6), "实验维度1: 精度验证", [
    "IV拟合：随机偏置 + 随机L/W/EOT组合",
    "CV拟合：QG, QS, QD vs VGS & VDS",
    "对比基准：校准后的BSIM-CMG模型卡",
    "指标：目视拟合质量 + RMS误差",
], tsz=17, bsz=14)
card(s, Inches(4.6), Inches(1.5), Inches(4), Inches(2.6), "实验维度2: 电路基准", [
    "1001级NAND振荡器（大规模数字）",
    "17级环形振荡器 + D触发器",
    "16位全加器（混合信号）",
    "NAND振荡器扫规模: 25K+1 ~ 100K+1级",
], tsz=17, bsz=14)
card(s, Inches(8.9), Inches(1.5), Inches(4), Inches(2.6), "实验维度3: 消融对比", [
    "激活函数: Sigmoid vs Tanh vs ISRU",
    "VA实现: 直接乘法 vs 循环矩阵乘法",
    "网络规模: 5 / 10 / 20 / 50 神经元",
    "对比基准: 商用BSIM-CMG Verilog-A",
], tsz=17, bsz=14)
card(s, Inches(0.3), Inches(4.4), Inches(12.3), Inches(2.5), "实验环境与评价指标 (Section III)", [
    "仿真器: SPICE (含Verilog-A编译)   |   数据源: 校准BSIM-CMG模型卡   |   温度: 27 C (常温)",
    "核心速度指标: 仿真总CPU时间 vs BSIM-CMG   |   收敛性指标: NR迭代次数对比   |   精度指标: 输出电压波形对齐",
    "测试逻辑: 先验证单管IV/CV精度 -> 再测试小规模电路(D触发器) -> 最后推到大规模电路(NAND振荡器 x 100K+级)",
], tsz=16, bsz=14)
notes(s.notes_slide.notes_text_frame, "实验设计分为三个维度。第一个维度是精度验证：在随机偏置和随机器件几何参数组合下，验证NN模型对IV和CV特性的拟合能力，对比基准是校准后的工业级BSIM-CMG模型。第二个维度是大规模电路基准测试：论文测试了从1001级到100K+1级不等的NAND振荡器，以及17级环形振荡器、D触发器和16位全加器等代表性电路。第三个维度是消融实验，系统对比了三种激活函数、两种VA实现方式（直接乘法vs循环）以及不同网络规模（5到50个神经元）下的仿真速度和收敛性差异。")

# ===== S8: 关键结果1 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(s)
hdr_bar(s, "NN模型在IV/CV拟合上达到BSIM-CMG级别精度，且仅需2x10=20个隐藏神经元", "#08 关键结果1: IV/CV拟合精度")
card(s, Inches(0.3), Inches(1.5), Inches(6.2), Inches(2.7), "IV拟合结果 (Fig. 2, Section II)", [
    "ID-VGS和ID-VDS曲线在不同L/HFIN/EOT下与BSIM-CMG几乎完全重合",
    "训练数据 vs 测试数据: 随机选取的偏置和几何参数组合",
    "2层x10神经元即可实现高精度拟合，无需更深网络",
    "亚阈值区、线性区、饱和区均准确捕获，无分区拼接痕迹",
], tsz=17, bsz=14)
card(s, Inches(6.8), Inches(1.5), Inches(6.2), Inches(2.7), "CV拟合结果 (Fig. 3, Section II)", [
    "C-V特性（QG,QS,QD）在VGS和VDS扫描下与BSIM-CMG高度一致",
    "关键：CV曲线光滑且导数连续，满足瞬态仿真的电荷守恒要求",
    "电荷偏移问题通过在损失函数中加入Q0项（VGS=VDS=0时的估计电荷）解决",
    "QV网络与IV网络共享输入特征但独立训练",
], tsz=17, bsz=14)
card(s, Inches(0.3), Inches(4.5), Inches(12.3), Inches(2.4), "精度达成路径总结", [
    "输出变换(1)(2)将ID/IG的宽动态范围（~12dec）映射到NN友好的对数空间 -> 加速收敛、提升数值稳定性",
    "导数辅助损失(3)使NN不仅学习电流值，还学习电导(gm,gds)及其斜率 -> SPICE仿真中Jacobian矩阵精确可计算",
    "QV网络直接输出终端电荷(非电容)，保证瞬态仿真的电荷守恒 -> 避免C(V)转Q(V)的数值积分误差",
], tsz=16, bsz=14)
notes(s.notes_slide.notes_text_frame, "第一个关键结果是IV和CV的拟合精度。论文的Fig.2展示了不同沟道长度L、鳍高HFIN和等效氧化层厚度EOT组合下的ID-VGS和ID-VDS曲线，NN模型的预测与BSIM-CMG几乎完全重合。令人印象深刻的是，只需2个隐藏层各10个神经元就能实现这一精度。Fig.3展示了CV特性的拟合结果，QG、QS、QD三个终端电荷在VGS和VDS扫描下均高度一致。这里需要强调的是，作者直接输出终端电荷而非电容，这保证了瞬态仿真中电荷守恒的数值精度。")

# ===== S9: 关键结果2 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(s)
hdr_bar(s, "直接乘法VA实现带来~30倍仿真加速，ISRU激活函数综合性能最优", "#09 关键结果2: 仿真速度与激活函数对比")
card(s, Inches(0.3), Inches(1.5), Inches(6.2), Inches(2.7), "仿真速度对比 (Table I, Fig. 7-8, Section III)", [
    "17级环振50ns仿真: BSIM-CMG 228s, ISRU直接乘法 7.8s (~29x加速)",
    "1001级NAND振荡器: NN模型与BSIM-CMG输出波形精确对齐，零收敛问题",
    "环形振荡器扫规模(25K~100K级): NN仿真时间曲线显著低于BSIM-CMG",
    "每次迭代平均耗时: NN模型约为BSIM-CMG的1/20 ~ 1/40",
], tsz=17, bsz=14)
card(s, Inches(6.8), Inches(1.5), Inches(6.2), Inches(2.7), "激活函数对比 (Table I, Section II)", [
    "ISRU: 仿真速度最快、NR迭代数最少（综合最优）",
    "Tanh: 精度与ISRU相当，但因含exp导致速度略慢",
    "Sigmoid: 速度和精度均落后于ISRU和Tanh",
    "选择原则: ISRU=f(x)=x/sqrt(1+x^2)，形式简洁且无指数运算",
], tsz=17, bsz=14)
card(s, Inches(0.3), Inches(4.5), Inches(12.3), Inches(2.4), "加速机制深度分析", [
    "直接乘法: 展开所有矩阵运算为 Wx+b 的显式表达式，VA编译器可逐行编译为C，避免解释型循环开销",
    "ISRU导数: d(ISRU)/dx = 1/(1+x^2)^(3/2)，在NR迭代中每步只需一次除法和开方，远快于tanh的exp/sinh/cosh",
    "与BSIM-CMG对比: BSIM-CMG每次模型评估需计算大量物理中间量(Vth, mobility, vsat...)，NN仅需20个神经元的乘加+激活",
], tsz=16, bsz=14)
notes(s.notes_slide.notes_text_frame, "第二个关键结果是仿真速度。在17级环形振荡器的50ns仿真中，BSIM-CMG耗时228秒，而ISRU NN模型直接乘法实现仅需7.8秒，加速约29倍。Fig.7展示了从25K+1到100K+1级NAND振荡器的仿真总时间，NN模型始终显著低于BSIM-CMG。激活函数的对比结果是ISRU综合最优，因为它没有指数运算，在SPICE的Newton-Raphson迭代中每步计算量最小。更深层次的原因是：BSIM-CMG每次模型评估需要计算大量的物理中间量，而NN仅需20个神经元的乘加运算加一次激活函数。")

# ===== S10: 消融/分析 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(s)
hdr_bar(s, "直接乘法 vs 循环实现的速度差异本质上是VA编译器优化的结果", "#10 消融实验与深入分析")
card(s, Inches(0.3), Inches(1.5), Inches(6.2), Inches(2.5), "VA实现方式对比 (Table I)", [
    "直接乘法: 展开矩阵乘为逐元素Wx+b表达式 -> 编译器逐行翻译为C",
    "循环实现: 使用for-loop + array -> 解释型执行，每步有循环开销",
    "在所有测试电路中直接乘法始终快于循环实现（差距随电路规模扩大）",
    "网络越大(50神经元)，直接乘法的加速优势越明显",
], tsz=17, bsz=14)
card(s, Inches(6.8), Inches(1.5), Inches(6.2), Inches(2.5), "网络规模与精度-速度权衡", [
    "5神经元: 速度最快但IV/CV精度不足以匹配BSIM-CMG",
    "10神经元: 精度达标且仿真速度显著快于BSIM-CMG（最佳平衡点）",
    "20神经元: 精度边际提升，但仿真速度开始下降",
    "50神经元: 精度饱和，仿真速度大幅下降（接近BSIM-CMG）",
], tsz=17, bsz=14)
card(s, Inches(0.3), Inches(4.3), Inches(12.3), Inches(2.6), "域外泛化与数值鲁棒性讨论", [
    "域外泛化: NN仅在训练数据范围内可保证精度，超出训练范围的偏置/几何组合可能导致不可预测的输出。物理模型天然具有外推能力",
    "数值稳定性: ISRU的导数在|x|大时趋近于1/|x|^2，不会完全消失（vs Tanh饱和区导数趋零），这在一定程度上缓解了NR求解器的Jacobian奇异问题",
    "局限性: 当前模型假设体浮空(Floating-body)，未考虑体电流IB；也未包含温度依赖性和工艺角建模",
], tsz=16, bsz=14)
notes(s.notes_slide.notes_text_frame, "本页的消融实验揭示了几个重要发现。第一，直接乘法实现相比循环实现的加速确实显著，原因在于VA编译器可以将展开的表达式逐行编译为C代码，避免了运行时的循环解释开销。第二，网络规模存在一个最优平衡点：10个神经元恰好可以提供足够的表达能力来匹配BSIM-CMG的精度，同时将仿真速度最大化。更多的神经元会导致精度饱和而速度下降。第三，我们需要坦然面对NN模型的一个固有局限：域外泛化能力。NN仅在训练数据覆盖的范围内保证精度，而物理模型天然具有外推能力。这是一个本质性的trade-off。")

# ===== S11: 创新点 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(s)
hdr_bar(s, "本文构建了首个经过完整电路级验证的SPICE兼容NN紧凑模型工具链", "#11 创新点与贡献总结")
card(s, Inches(0.3), Inches(1.5), Inches(4), Inches(3.5), "方法创新", [
    "ISRU激活函数的首次系统应用于NN紧凑模型场景",
    "导数辅助损失（含二阶导数）用于保证SPICE收敛性",
    "输出变换公式 ID=VDS*exp(y1) 将物理约束嵌入NN",
    "直接乘法VA实现策略的系统对比验证",
], tsz=17, bsz=14)
card(s, Inches(4.6), Inches(1.5), Inches(4), Inches(3.5), "实验创新", [
    "覆盖100K+级NAND振荡器的大规模电路验证",
    "多种电路类型（数字/混合信号/振荡器）的基准测试",
    "三种激活函数 + 两种VA实现 + 四种网络规模的系统消融",
    "与工业级BSIM-CMG的端到端精度+速度对比",
], tsz=17, bsz=14)
card(s, Inches(8.9), Inches(1.5), Inches(4), Inches(3.5), "工具链贡献", [
    "Python自动代码生成器: 训练好的权重直接导出编译就绪的Verilog-A",
    "训练->VA导出->OpenVAF编译->SPICE仿真的全自动化流程",
    "为AI/ML紧凑模型社区提供了可复现的基准和最佳实践",
], tsz=17, bsz=14)
rrect(s, Inches(0.3), Inches(5.3), Inches(12.3), Inches(1.6), C_DARK)
rect(s, Inches(0.35), Inches(5.35), Inches(0.05), Inches(1.5), C_ACC2)
ml(s, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.3), [
    "论文核心贡献一句话总结: 首次证明了2层10神经元ISRU-NN可以在大规模电路仿真中以~1/30的CPU时间替代BSIM-CMG，且无收敛问题",
    "* 该成果已被延伸为更完整的BSIM-NN框架（TED 2023, 2025），包含NQS、自热、噪声、统计变异等高级物理效应",
], sz=16, c=C_WHITE)
notes(s.notes_slide.notes_text_frame, "总结本文的创新和贡献。方法层面：ISRU激活函数的选择、导数辅助损失函数的设计、输出变换的物理约束嵌入以及直接乘法的VA实现策略，构成了一个完整的SPICE兼容NN模型方法论。实验层面：覆盖了从单管到100K+级大规模电路的完整基准测试，这在之前的NN紧凑模型工作中是前所未有的。工具链层面：作者开发的Python自动VA代码生成器打通了从训练到仿真的最后一公里。这项工作为后续的BSIM-NN完整框架奠定了基础，目前该框架已扩展到包含NQS、自热、噪声和统计变异等高级效应。")

# ===== S12: 不足与讨论 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(s)
hdr_bar(s, "NN紧凑模型的工业级落地仍需解决域外泛化、实测数据验证和多物理效应覆盖等问题", "#12 不足、启发与讨论")
card(s, Inches(0.3), Inches(1.5), Inches(6), Inches(2.4), "本文的不足与局限", [
    "仅使用BSIM-CMG合成数据训练，未在硅实测数据上验证（噪声、局部变异鲁棒性未知）",
    "域外泛化无保证: 超出训练范围的偏置/几何条件下NN输出不可控",
    "假设体浮空、无温度依赖性（2023/2025版已改进为包含SH/NQS/温度）",
    "未讨论NN模型的重新训练成本是否低于传统参数提取",
], tsz=17, bsz=14)
card(s, Inches(6.6), Inches(1.5), Inches(6), Inches(2.4), "对自身工作的启发", [
    "ISRU激活 + 导数损失的设计可直接应用于我们的3-32-32-1模型",
    "输出变换 ID=VDS*exp(y1) 优于我们当前的 log10(ID) 方案（天然保证VDS=0边界条件）",
    "直接乘法VA实现是我们当前代码中已采用的最佳实践",
    "McAndrew基准清单为我们的模型验证提供了标准化测试框架",
], tsz=17, bsz=14)
card(s, Inches(0.3), Inches(4.2), Inches(12.3), Inches(2.6), "开放讨论问题", [
    "1. NN紧凑模型的"可解释性"是否真的重要？还是只要电路仿真结果正确即可？",
    "2. 当工艺节点迁移时（如N5->N3），NN模型是否需要完全重新训练？能否用迁移学习降低数据成本？",
    "3. 对于混合信号设计中的噪声/匹配分析，NN模型能否提供与物理模型同等的统计仿真能力？",
    "4. 开放科学问题: NN紧凑模型的收敛性是否可以在理论上给出保证（而不仅仅通过实验验证）？",
], tsz=17, bsz=15)
notes(s.notes_slide.notes_text_frame, "最后一页我们讨论本文的不足和对我们自身工作的启发。本文的主要局限包括仅使用合成数据验证、域外泛化能力未讨论、以及某些物理效应的简化假设。这些在后续的BSIM-NN工作中已有改进。对我们的项目而言，ISRU激活函数、导数辅助损失和输出变换设计是可以直接借鉴的改进方向。最后我提出四个开放讨论问题：NN模型可解释性的必要性、工艺迁移时的重训练成本、统计仿真的能力、以及NN收敛性的理论保证——这些是当前领域内尚未完全解决的深层问题，值得我们在后续研究中思考。")

prs.save(out)
print("PPT saved: %s" % out)
print("Slides: %d" % len(prs.slides))