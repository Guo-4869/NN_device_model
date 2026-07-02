import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

outpath = r"D:\gzp\研究生\半导体器件建模与仿真\NN_device_model\项目汇报.pptx"
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

C_DARK = RGBColor(0x1a, 0x1a, 0x2e)
C_ACC = RGBColor(0x00, 0x6d, 0x77)
C_ACC2 = RGBColor(0xe8, 0x95, 0x22)
C_LIGHT = RGBColor(0xf0, 0xf0, 0xf5)
C_W = RGBColor(0xff, 0xff, 0xff)
C_GRAY = RGBColor(0x66, 0x66, 0x66)
C_RED = RGBColor(0xcc, 0x33, 0x33)
C_GREEN = RGBColor(0x2e, 0x7d, 0x32)
C_CARD = RGBColor(0xef, 0xef, 0xf2)
C_SUB = RGBColor(0xaa, 0xaa, 0xaa)

def bg(s, c=C_DARK):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = c

def shp(s, l, t, w, h, c):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = c
    sh.line.fill.background()
    return sh

def rshp(s, l, t, w, h, c):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = c
    sh.line.fill.background()
    return sh

def tb(s, l, t, w, h, txt, sz=18, b=False, c=C_W, a=PP_ALIGN.LEFT):
    tb_ = s.shapes.add_textbox(l, t, w, h)
    tf = tb_.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = txt
    p.font.size = Pt(sz); p.font.bold = b
    p.font.color.rgb = c; p.font.name = 'Microsoft YaHei'; p.alignment = a

def ml(s, l, t, w, h, lines, sz=16, c=C_W, ls=1.3):
    tb_ = s.shapes.add_textbox(l, t, w, h)
    tf = tb_.text_frame; tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln; p.font.size = Pt(sz)
        p.font.color.rgb = c; p.font.name = 'Microsoft YaHei'
        p.space_after = Pt(sz * (ls - 1))

def hdr(s, t, sub=""):
    bg(s, C_DARK)
    shp(s, Inches(0), Inches(3.25), prs.slide_width, Inches(1), C_ACC)
    tb(s, Inches(1), Inches(3.35), Inches(11), Inches(0.8), t, sz=36, b=True)
    if sub:
        tb(s, Inches(1), Inches(4.2), Inches(11), Inches(0.6), sub, sz=16, c=C_SUB)

def card(s, l, t, w, h, title, bullets, tsz=18, bsz=14):
    rshp(s, l, t, w, h, C_CARD)
    y = Emu(t.emu + Inches(0.15).emu)
    tb(s, Emu(l.emu + Inches(0.2).emu), y, Emu(w.emu - Inches(0.4).emu), Inches(0.4), title, sz=tsz, b=True, c=C_DARK)
    y = Emu(t.emu + Inches(0.6).emu)
    lines = ["\u2022  " + b for b in bullets]
    ml(s, Emu(l.emu + Inches(0.2).emu), y, Emu(w.emu - Inches(0.4).emu), Emu(h.emu - Inches(0.7).emu), lines, sz=bsz, c=C_GRAY)

# ===== Slide 1: Title =====
s1 = prs.slides.add_slide(prs.slide_layouts[6])
bg(s1, C_DARK)
shp(s1, Inches(0), Inches(0), prs.slide_width, Inches(0.06), C_ACC)
tb(s1, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2), "基于神经网络的半导体器件紧凑建模", sz=44, b=True)
tb(s1, Inches(1.5), Inches(3.1), Inches(10), Inches(0.7), "BSIM-NN: Neural Network-Based MOSFET Compact Model", sz=22, c=C_ACC)
shp(s1, Inches(1.5), Inches(3.85), Inches(2.5), Inches(0.04), C_ACC2)
tb(s1, Inches(1.5), Inches(4.1), Inches(10), Inches(0.5), "从 PyTorch 训练到 ngspice 电路仿真的端到端流程", sz=18, c=C_SUB)
tb(s1, Inches(1.5), Inches(4.6), Inches(10), Inches(0.5), "汇报人：Guo-4869  |  半导体器件建模与仿真", sz=14, c=C_SUB)
tb(s1, Inches(1.5), Inches(5.0), Inches(10), Inches(0.5), "基于开源工具：PyTorch + OpenVAF + ngspice", sz=14, c=C_SUB)

# ===== Slide 2: 背景 =====
s2 = prs.slides.add_slide(prs.slide_layouts[6])
hdr(s2, "项目背景与动机")
card(s2, Inches(0.5), Inches(0.8), Inches(5.8), Inches(2.6), "传统BSIM紧凑模型的局限", [
    "基于物理方程的解析公式，开发周期长",
    "新器件（GAA、Nanosheet）建模需数年",
    "参数提取过程复杂，需专业工程师经验",
    "模型精度与仿真速度难以同时优化",
], tsz=18, bsz=14)
card(s2, Inches(6.8), Inches(0.8), Inches(5.8), Inches(2.6), "NN紧凑模型的优势", [
    "数据驱动替代物理推导，缩短开发周期",
    "从IV数据自动学习器件行为，捕获非线性",
    "可拟合任意器件特性，不依赖特定物理假设",
    "兼容标准SPICE仿真流程（VA/OSDI）",
], tsz=18, bsz=14)
card(s2, Inches(0.5), Inches(3.8), Inches(12.3), Inches(3.1), "核心技术路线", [
    "训练数据：BSIM3解析模型批量计算IV特性（VGS/VDS/VBS扫描）",
    "NN训练：PyTorch MLP拟合log10(ID)，含权重正则化和体偏置",
    "模型导出：PyTorch权重 -> Verilog-A代码，tanh激活实现光滑可微",
    "编译仿真：OpenVAF编译VA -> OSDI共享库 -> ngspice电路级仿真",
    "核心挑战：NN导数光滑性直接影响SPICE Newton-Raphson迭代收敛",
], tsz=18, bsz=14)

# ===== Slide 3: 路线图 =====
s3 = prs.slides.add_slide(prs.slide_layouts[6])
hdr(s3, "技术路线与工作流程")
steps = [
    ("1. 数据生成", ["BSIM3解析模型", "VGS/VDS/VBS扫描", "-> 12500样本"]),
    ("2. PyTorch训练", ["3->32->32->1 MLP", "Tanh+AdamW", "-> .pth模型"]),
    ("3. VA导出", ["权重自动转换", "smooth tanh钳位", "-> .va文件"]),
    ("4. OpenVAF编译", ["VA->OSDI共享库", "ngspice加载", "-> .osdi文件"]),
    ("5. 电路仿真", ["DC/AC/瞬态", "电路级验证", "-> 模型即用"]),
]
for i, (ttl, buls) in enumerate(steps):
    x = Inches(0.3 + i * 2.6)
    card(s3, x, Inches(1.2), Inches(2.3), Inches(2.8), ttl, buls, tsz=16, bsz=13)
    if i < 4:
        shp(s3, Inches(0.3 + i * 2.6 + 2.4), Inches(2.4), Inches(0.15), Inches(0.04), C_ACC2)
tb(s3, Inches(0.5), Inches(4.5), Inches(12), Inches(0.5), "技术栈", sz=20, b=True, c=C_W)
techs = ["Python", "NumPy", "PyTorch", "scikit-learn", "Matplotlib", "OpenVAF", "ngspice"]
for i, tt in enumerate(techs):
    card(s3, Inches(0.5 + i * 1.8), Inches(5.1), Inches(1.6), Inches(1.2), "", [tt], tsz=16, bsz=14)

# ===== Slide 4: NN架构 =====
s4 = prs.slides.add_slide(prs.slide_layouts[6])
hdr(s4, "神经网络模型架构")
card(s4, Inches(0.5), Inches(1.0), Inches(4), Inches(1.2), "输入层（3维）", ["VGS 栅源电压", "VDS 漏源电压", "VBS 体偏置电压"], tsz=18, bsz=14)
card(s4, Inches(4.8), Inches(1.0), Inches(4), Inches(1.2), "隐藏层1（32神经元）", ["全连接 fc1: 3->32", "激活: tanh(x)", "输出范围: (-1,1)"], tsz=18, bsz=14)
card(s4, Inches(9.1), Inches(1.0), Inches(4), Inches(1.2), "隐藏层2（32神经元）", ["全连接 fc2: 32->32", "激活: tanh(x)", "权重 < 5.5"], tsz=18, bsz=14)
card(s4, Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.5), "输出层", [
    "全连接 fc3: 32->1（线性），输出 log10(ID)",
    "ID = 10^log10(ID) x (W/10um)  +  光滑VDSe0过渡 + 输出钳位",
], tsz=18, bsz=14)
card(s4, Inches(0.5), Inches(4.3), Inches(6), Inches(2.6), "设计关键", [
    "tanh激活: 导数=1-tanh^2(x), VA内置函数, 无额外指数",
    "log10(ID)输出: 对数空间线性关系, 覆盖 -15到-3 (12dec)",
    "L2正则化: weight_decay=1e-4, 限制权重防tanh饱和",
    "VBS训练: 包含VBS扫描, 保证归一化std>0.05",
], tsz=18, bsz=14)
card(s4, Inches(6.8), Inches(4.3), Inches(6), Inches(2.6), "模型统计", [
    "总参数: 1,217 (3x32 + 32x32 + 32x1 + bias)",
    "训练样本: 12,500 (50x50 VGS/VDS x 5 VBS)",
    "优化器: AdamW (lr=0.001, wd=1e-4)",
    "损失: MSE on log10(ID) + ReduceLROnPlateau",
], tsz=18, bsz=14)

# ===== Slide 5: 版本演进 =====
s5 = prs.slides.add_slide(prs.slide_layouts[6])
hdr(s5, "版本演进历程", "从概念验证到SPICE兼容的迭代开发")
vers = [
    ("V1", "概念验证", "2in->10 Tanh\n简化MOSFET数据", RGBColor(0xcc,0x33,0x33)),
    ("V2", "加权训练", "亚阈值区域加权\n改善弱反型拟合", RGBColor(0xe8,0x6a,0x1a)),
    ("V3", "多输入MLP", "5in->16->16->3\nBSIM3数据生成器", RGBColor(0xe8,0x95,0x22)),
    ("V4", "论文实现", "ISRU激活\ngm/gds导数损失", RGBColor(0x2e,0x7d,0x32)),
    ("V5", "简化收敛版", "3->16->16->1\nOSDI仿真卡死", RGBColor(0x00,0x6d,0x77)),
    ("V6", "SPICE兼容", "3->32->32->1\n全修复,不卡死", C_ACC),
]
for i, (v, ttl, desc, clr) in enumerate(vers):
    x = Inches(0.3 + i * 2.15)
    rshp(s5, x, Inches(1.3), Inches(2.0), Inches(3.0), clr)
    tb(s5, x + Inches(0.15), Inches(1.4), Inches(1.7), Inches(0.5), v, sz=24, b=True)
    tb(s5, x + Inches(0.15), Inches(1.85), Inches(1.7), Inches(0.4), ttl, sz=16, b=True)
    ml(s5, x + Inches(0.15), Inches(2.3), Inches(1.7), Inches(1.8), desc.split("\n"), sz=12)

card(s5, Inches(0.5), Inches(4.6), Inches(12.3), Inches(2.3), "V5 -> V6: OSDI仿真卡死关键修复", [
    "V5: if(VDS==0) -> V6: 有理函数  Ids*VDS/(VDS+1e-4)  光滑过渡",
    "V5: std_vbs=0.000001 -> V6: VBS扫描训练, 强制std>=0.05",
    "V5: 二层权重-7.66 -> V6: L2正则化 所有权重<5.5",
    "V5: ?:三元钳位 -> V6: tanh软钳位  log10Id=-7.5+4.5*tanh((x+7.5)/4.5)",
    "V5: I(g,s)<+0 -> V6: 小信号电导 I(g,s)<+1e-14*V(g,s)",
], tsz=18, bsz=13)

# ===== Slide 6: 训练结果 =====
s6 = prs.slides.add_slide(prs.slide_layouts[6])
hdr(s6, "Version 6 训练结果")
card(s6, Inches(0.5), Inches(1.0), Inches(5.8), Inches(3.0), "训练精度指标", [
    "验证MSE: 0.000410",
    "log10(ID) MAE: 0.0140",
    "log10(ID) RMSE: 0.0228",
    "ID中位相对误差: 2.00%",
    "ID P95相对误差: 10.08%",
    "训练轮数: 1,500 epochs",
], tsz=18, bsz=14)
card(s6, Inches(6.8), Inches(1.0), Inches(5.8), Inches(3.0), "权重控制（SPICE收敛关键）", [
    "fc1.weight: max=2.74 [OK]",
    "fc2.weight: max=5.50 [略高]",
    "fc3.weight: max=1.89 [OK]",
    "所有bias max<1.31 [OK]",
    "vs V5二层-7.66 [改善40%+]",
    "训练收敛稳定，无过拟合",
], tsz=18, bsz=14)
img = r"D:\gzp\研究生\半导体器件建模与仿真\NN_device_model\version6\v6_results.png"
if os.path.exists(img):
    s6.shapes.add_picture(img, Inches(0.5), Inches(4.2), Inches(12.3), Inches(3.0))

# ===== Slide 7: VA =====
s7 = prs.slides.add_slide(prs.slide_layouts[6])
hdr(s7, "Verilog-A 模型导出与编译")
card(s7, Inches(0.5), Inches(1.0), Inches(6.0), Inches(2.8), "VA代码生成（133行）", [
    "自动提取权重 -> 完整 .va 源文件",
    "3in->32tanh->32tanh->1linear->log10(ID)",
    "光滑输入: VGS=0.5+0.7*tanh((x-0.5)/0.7)",
    "光滑输出: log10Id=-7.5+4.5*tanh((x+7.5)/4.5)",
    "电流: Ids=10^log10Id x (W/10um)",
], tsz=18, bsz=14)
card(s7, Inches(6.8), Inches(1.0), Inches(6.0), Inches(2.8), "编译与仿真", [
    "openvaf bsim_nn_v6.va --ngspice -o ...osdi",
    "ngspice: .control pre_osdi ...osdi",
    ".model NMOS1 bsim_nn_v6(w=10u l=1u)",
    "支持 DC/OP/AC 分析",
    "附 test_v6.cir 测试电路",
], tsz=18, bsz=14)
card(s7, Inches(0.5), Inches(4.1), Inches(12.3), Inches(2.8), "SPICE兼容性保障", [
    "所有激活用VA内置tanh，OpenVAF可正确求导，无需自定义函数",
    "光滑VDS->0过渡避免if导致的导数不连续，助力NR收敛",
    "Gate/Body添加1e-14*V小信号电导，保证Jacobian非奇异",
    "归一化std>=0.05杜绝除零溢出，避免求解器崩溃",
    "权重正则化限制tanh输入范围，防饱和导致Jacobian病态",
], tsz=18, bsz=14)

# ===== Slide 8: 总结 =====
s8 = prs.slides.add_slide(prs.slide_layouts[6])
hdr(s8, "总结与展望")
card(s8, Inches(0.5), Inches(1.0), Inches(5.8), Inches(3.0), "当前成果", [
    "完成 PyTorch -> ngspice 完整工具链",
    "针对SPICE收敛实施5项关键修复",
    "log10(ID) MAE=0.014, ID误差中位2%",
    "权重控制良好 (max<5.5)",
    "已开源: github.com/Guo-4869/NN_device_model",
], tsz=18, bsz=14)
card(s8, Inches(6.8), Inches(1.0), Inches(5.8), Inches(3.0), "下一步", [
    "用实际BSIM4模型批量仿真生成数据",
    "拓展3层网络 (32->32->32) 提升精度",
    "加入gm/gds导数辅助损失",
    "扩展输入: 沟长L/宽度W/温度T",
    "测试实际电路: 反相器/差分对等",
], tsz=18, bsz=14)
card(s8, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.5), "未来方向", [
    "更复杂NN架构: ResNet/Attention的SPICE兼容实现",
    "多器件联合建模: 同时拟合NMOS+PMOS保证电路仿真一致性",
    "自动化模型生成: 从工艺PDK自动生成NN紧凑模型，减少人工建模",
    "开源贡献: 为OpenVAF/ngspice社区提供NN紧凑模型参考实现",
], tsz=18, bsz=14)

prs.save(outpath)
print(f"PPT saved: {outpath}")
print(f"Slides: {len(prs.slides)}")